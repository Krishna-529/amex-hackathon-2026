import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Palette, matched to assets/deck/PPT.pptx (the team's reference deck) ----
    C_BG = RGBColor(0x0A, 0x1B, 0x3D)          # slide background — deep navy
    C_BG_GLOW1 = RGBColor(0x0E, 0x24, 0x50)    # decorative glow layers on the title slide
    C_BG_GLOW2 = RGBColor(0x10, 0x27, 0x5A)
    C_BG_GLOW3 = RGBColor(0x16, 0x33, 0x6F)
    C_CARD = RGBColor(0x13, 0x2C, 0x52)        # panel / card fill
    C_CARD_DIM = RGBColor(0x0D, 0x22, 0x46)    # header rows, small chips
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    C_BODY = RGBColor(0xA9, 0xBE, 0xDC)        # body copy on navy
    C_MUTED = RGBColor(0x7B, 0x94, 0xBA)       # captions / meta / source lines
    C_ACCENT = RGBColor(0x3D, 0x9C, 0xFF)      # electric blue — the deck's signature accent
    C_GOOD = RGBColor(0x2E, 0xCC, 0x9A)        # green
    C_BAD = RGBColor(0xFF, 0x6B, 0x6B)         # coral red
    C_WARN = RGBColor(0xFF, 0xB0, 0x20)        # amber
    C_PURPLE = RGBColor(0x9B, 0x8C, 0xFF)      # violet

    FONT_HEAD = "Cambria"
    FONT_BODY = "Calibri"
    FONT_MONO = "Courier New"

    blank_layout = prs.slide_layouts[6]

    # ---------------- helpers ----------------

    def set_background(slide, color=C_BG):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def no_shadow(shape):
        shape.shadow.inherit = False

    def panel(slide, left, top, width, height, fill=C_CARD, radius=0.06):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        try:
            card.adjustments[0] = radius
        except Exception:
            pass
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.fill.background()
        no_shadow(card)
        return card

    def pad(text_frame, inset=0.28):
        text_frame.margin_left = Inches(inset)
        text_frame.margin_right = Inches(inset)
        text_frame.margin_top = Inches(inset)
        text_frame.margin_bottom = Inches(inset)
        text_frame.word_wrap = True

    def add_header(slide, section_no, section_name, title_text, subtitle_text, page_no):
        set_background(slide)

        kicker_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.16), Inches(9.0), Inches(0.28))
        ktf = kicker_box.text_frame
        ktf.word_wrap = True
        kp = ktf.paragraphs[0]
        kp.text = f"{section_no} · {section_name.upper()}"
        kp.font.size = Pt(11.5)
        kp.font.bold = True
        kp.font.name = FONT_BODY
        kp.font.color.rgb = C_ACCENT

        pg_box = slide.shapes.add_textbox(Inches(12.2), Inches(0.16), Inches(0.6), Inches(0.28))
        ptf = pg_box.text_frame
        pp = ptf.paragraphs[0]
        pp.text = str(page_no)
        pp.font.size = Pt(11.5)
        pp.font.name = FONT_BODY
        pp.font.color.rgb = C_MUTED
        pp.alignment = PP_ALIGN.RIGHT

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11.9), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(25)
        p.font.bold = True
        p.font.name = FONT_HEAD
        p.font.color.rgb = C_WHITE

        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.02), Inches(11.9), Inches(0.4))
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = subtitle_text
            sp.font.size = Pt(13)
            sp.font.name = FONT_BODY
            sp.font.color.rgb = C_BODY

    def panel_label(text_frame, label, color, title, space_after=10):
        p = text_frame.paragraphs[0]
        p.text = label.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = FONT_BODY
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p2 = text_frame.add_paragraph()
        p2.text = title
        p2.font.size = Pt(16.5)
        p2.font.bold = True
        p2.font.name = FONT_HEAD
        p2.font.color.rgb = C_WHITE
        p2.space_after = Pt(space_after)

    def bullet_block(text_frame, items, size=12.5, marker_color=C_ACCENT, space_after=8, first_new=False):
        for i, b in enumerate(items):
            p = text_frame.paragraphs[0] if (i == 0 and not first_new) else text_frame.add_paragraph()
            p.text = "—  " + b
            p.font.size = Pt(size)
            p.font.name = FONT_BODY
            p.font.color.rgb = C_BODY
            p.space_after = Pt(space_after)
            p.line_spacing = 1.2

    def source_line(slide, text, top=Inches(7.08)):
        box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.name = FONT_BODY
        p.font.color.rgb = C_MUTED

    # ==================== SLIDE 1: Title ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)

    for (color, left, top, size, rot) in [
        (C_BG_GLOW1, Inches(-3.2), Inches(-2.6), Inches(8.4), 12),
        (C_BG_GLOW2, Inches(9.1), Inches(-1.9), Inches(6.6), -8),
        (C_BG_GLOW3, Inches(10.4), Inches(-0.6), Inches(3.9), 20),
    ]:
        blob = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size)
        blob.rotation = rot
        try:
            blob.adjustments[0] = 0.5
        except Exception:
            pass
        blob.fill.solid()
        blob.fill.fore_color.rgb = color
        blob.line.fill.background()
        no_shadow(blob)

    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.0), Inches(2.7))
    tf1 = tb1.text_frame
    pad(tf1, 0)

    p0 = tf1.paragraphs[0]
    p0.text = "CODESTREET 2026  ·  AMERICAN EXPRESS"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.name = FONT_BODY
    p0.font.color.rgb = C_ACCENT
    p0.space_after = Pt(16)

    p1 = tf1.add_paragraph()
    p1.text = "ZKD Concierge"
    p1.font.size = Pt(50)
    p1.font.bold = True
    p1.font.name = FONT_HEAD
    p1.font.color.rgb = C_WHITE
    p1.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "An autonomous travel-disruption concierge for Indian domestic aviation"
    p2.font.size = Pt(19)
    p2.font.name = FONT_BODY
    p2.font.color.rgb = C_BODY
    p2.space_after = Pt(4)

    stats = [
        ("~11 s", "cancellation to a confirmed,\npaid-for trip"),
        ("42 s", "of the work already done\nbefore it happens"),
        ("3", "consent modes, one\ndefault-deny policy gate"),
        ("100%", "agent actions passed\nthrough OPA / Rego"),
    ]
    for i, (num, desc) in enumerate(stats):
        left = Inches(1.0 + i * 2.86)
        card = panel(slide1, left, Inches(5.05), Inches(2.66), Inches(1.35), fill=C_CARD_DIM)
        tf = card.text_frame
        pad(tf, 0.2)
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(25)
        p.font.bold = True
        p.font.name = FONT_HEAD
        p.font.color.rgb = C_ACCENT
        p.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.text = desc.replace("\n", " ")
        p2.font.size = Pt(10.5)
        p2.font.name = FONT_BODY
        p2.font.color.rgb = C_BODY
        p2.line_spacing = 1.1

    footer = slide1.shapes.add_textbox(Inches(1.0), Inches(6.75), Inches(11.0), Inches(0.4))
    ftf = footer.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Team ZKD · IIT Madras   —   Live MVP presentation (20 min pitch + 10 min Q&A)"
    fp.font.size = Pt(12)
    fp.font.name = FONT_BODY
    fp.font.color.rgb = C_MUTED

    # ==================== SLIDE 2: The Passenger Reality ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "01", "The Reality", "The passenger reality: from panic to peace of mind",
                "What changes for the member the moment a flight is cancelled.", 2)

    box_old = panel(slide2, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.15))
    tf_old = box_old.text_frame
    pad(tf_old)
    panel_label(tf_old, "TODAY", C_BAD, "The reactive IRROPS nightmare")
    bullet_block(tf_old, [
        "Discovering the cancellation only at the gate or via a delayed SMS.",
        "45+ minute hold times with airline customer care centres.",
        "Scrambling for overpriced alternate flights and airport hotels.",
        "Opaque DGCA duty-of-care claims requiring manual paperwork.",
    ], marker_color=C_BAD, first_new=True)

    box_new = panel(slide2, Inches(6.65), Inches(1.55), Inches(6.18), Inches(5.15))
    tf_new = box_new.text_frame
    pad(tf_new)
    panel_label(tf_new, "ZKD Concierge", C_GOOD, "Proactive and autonomous")
    bullet_block(tf_new, [
        "Predictive warning ~42s before the formal airline cancellation.",
        "Instant push notification with pre-cached flight, hotel and ground options.",
        "1-tap pre-authorisation, or autopilot protection, with no phone queues.",
        "Automatic compensation claims and zero out-of-pocket stress.",
    ], marker_color=C_GOOD, first_new=True)

    # ==================== SLIDE 3: Live MVP Demo Flow ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "02", "Live Demo", "Live MVP demo walkthrough",
                "The pre-tea-break hero flow, run end to end on the live app.", 3)

    steps = [
        ("01", C_ACCENT, "Operator console (/ops)", "Trigger test disruptions locally or hosted; exercises the webhook and poller simulation channels instantly."),
        ("02", C_PURPLE, "Risk & LLM explanation", "Live risk gauge crossing thresholds, with a plain-language Gemini explanation of the weather/airline risk."),
        ("03", C_GOOD, "Edge case: self-cancel", "Member-initiated changes, a clean hand-off (handed-over), and zero-friction state sync."),
        ("04", C_WARN, "11-second autonomous recovery", "Pre-auth consent gate, multi-supplier option ranking, and 1-tap booking confirmation in ~11s."),
    ]

    for i, (num, color, stitle, sdesc) in enumerate(steps):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.18)
        top = Inches(1.55 + row * 2.55)

        card = panel(slide3, left, top, Inches(5.9), Inches(2.3))
        tf = card.text_frame
        pad(tf)

        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.name = FONT_MONO
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p_t = tf.add_paragraph()
        p_t.text = stitle
        p_t.font.size = Pt(15.5)
        p_t.font.bold = True
        p_t.font.name = FONT_HEAD
        p_t.font.color.rgb = C_WHITE
        p_t.space_after = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12.5)
        p2.font.name = FONT_BODY
        p2.font.color.rgb = C_BODY
        p2.line_spacing = 1.2

    # ==================== SLIDE 4: Predictive ML + Personalization ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "03", "Prediction & Personalization", "Predictive ML model & personalization",
                "The Spotify analogy: the same curation instinct, applied to a cancelled flight.", 4)

    box_ml1 = panel(slide4, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.15))
    tf1 = box_ml1.text_frame
    pad(tf1)
    panel_label(tf1, "Risk model", C_ACCENT, "Self-trained cancellation risk model")
    bullet_block(tf1, [
        "Built in zkd-risk-model/: XGBoost trained on real US DOT/BTS and Brazil ANAC historical data.",
        "Extracts weather, route congestion, and airline rolling delays.",
        "Outputs a continuous probability score with adaptive thresholds — not a flat 25/55/80.",
        "Advisory role only: decides when we start preparing, never whether we spend member money.",
    ], first_new=True)

    box_ml2 = panel(slide4, Inches(6.65), Inches(1.55), Inches(6.18), Inches(5.15))
    tf2 = box_ml2.text_frame
    pad(tf2)
    panel_label(tf2, "Personalization", C_PURPLE, "The Spotify / YouTube analogy")
    bullet_block(tf2, [
        "Just as Spotify curates Discover Weekly from listening history, ZKD learns cardmember preferences.",
        "Captures seat preference, cabin entitlement, layover tolerance, and companion rules.",
        "Translates free-text member intent (server/preferences/intent.ts) into weighted scoring rules.",
        "Dynamic adaptation: refines options from implicit choices and explicit feedback.",
    ], first_new=True)

    # ==================== SLIDE 5: Lifecycle Workflow ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "04", "The Lifecycle", "The disruption lifecycle",
                "Prediction buys speed: every phase before consent runs before the airline even confirms.", 5)

    phases = [
        ("01", C_ACCENT, "Predict", "~42s head start", "Model forecasts cancellation probability before the airline files it. Alt flights and hotels are pre-cached warm."),
        ("02", C_PURPLE, "Pre-cache", "Zero spend, no hold", "Options are priced and held warm in memory. No booking, no spend, no financial liability."),
        ("03", C_GOOD, "Consent gate", "Member control", "A high-risk threshold triggers pre-auth. The member chooses, or approves autopilot in advance."),
        ("04", C_WARN, "Act & settle", "~11s recovery", "The confirmed booking executes instantly. DGCA duty-of-care claims are filed automatically."),
    ]

    for i, (num, color, ptitle, psub, pdesc) in enumerate(phases):
        left = Inches(0.5 + i * 3.06)
        card = panel(slide5, left, Inches(1.65), Inches(2.86), Inches(4.9))
        tf = card.text_frame
        pad(tf, 0.24)

        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.name = FONT_MONO
        p.font.color.rgb = color
        p.space_after = Pt(6)

        p_t = tf.add_paragraph()
        p_t.text = ptitle
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.name = FONT_HEAD
        p_t.font.color.rgb = C_WHITE
        p_t.space_after = Pt(2)

        p_sub = tf.add_paragraph()
        p_sub.text = psub
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.name = FONT_BODY
        p_sub.font.color.rgb = color
        p_sub.space_after = Pt(14)

        p_desc = tf.add_paragraph()
        p_desc.text = pdesc
        p_desc.font.size = Pt(11.5)
        p_desc.font.name = FONT_BODY
        p_desc.font.color.rgb = C_BODY
        p_desc.line_spacing = 1.25

    # ==================== SLIDE 6: Policy Intelligence via RAG ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "05", "Policy & Compliance", "Policy intelligence via RAG",
                "Multi-regulation compliance, enforced by a default-deny policy gate.", 6)

    box_rag1 = panel(slide6, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.15))
    tf1 = box_rag1.text_frame
    pad(tf1)
    panel_label(tf1, "Scope", C_ACCENT, "Multi-regulation scope (beyond DGCA)")
    bullet_block(tf1, [
        "Not just DGCA (Indian domestic aviation) — international frameworks too: EU261 and UK261.",
        "Carrier-specific contract clauses and operating terms (e.g. British Airways vs Emirates entitlements).",
        "Amex card product benefit terms (server/myca.ts): cabin ceilings, lounge access, per-transaction rules.",
        "Jurisdiction routing automatically detects which regulatory framework governs the PNR.",
    ], first_new=True)

    box_rag2 = panel(slide6, Inches(6.65), Inches(1.55), Inches(6.18), Inches(5.15))
    tf2 = box_rag2.text_frame
    pad(tf2)
    panel_label(tf2, "Engine", C_WARN, "RAG & default-deny policy engine")
    bullet_block(tf2, [
        "RAG ingestion automatically vectorizes and indexes aviation regulations, fare rules, and card terms.",
        "Context retrieval dynamically queries relevant clauses when a disruption occurs for a given PNR.",
        "OPA policy gate (server/policy/): default-deny rule engine, evaluated in-process in under 1ms.",
        "Incomplete policy inputs trigger default-deny — nothing executes without an explicit allow.",
    ], first_new=True)

    # ==================== SLIDE 7: Architecture & Multi-Device Sync ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "06", "Architecture", "System architecture & multi-device sync",
                "A server-authoritative engine, and clients that are honest about holding no state.", 7)

    box_arch1 = panel(slide7, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.15))
    tf1 = box_arch1.text_frame
    pad(tf1)
    panel_label(tf1, "Engine", C_ACCENT, "Server-authoritative simulation engine")
    bullet_block(tf1, [
        "Module-level simulation engine (server/engine/simulation.ts) runs the whole lifecycle.",
        "Real setTimeout/setInterval chains resolve on schedule whether devices watch or not.",
        "Single source of truth: the shared state ledger (server/decisionLedger.ts).",
        "Zero GPU on the critical path — 95% of the time is network wait on supplier APIs.",
    ], first_new=True)

    box_arch2 = panel(slide7, Inches(6.65), Inches(1.55), Inches(6.18), Inches(5.15))
    tf2 = box_arch2.text_frame
    pad(tf2)
    panel_label(tf2, "Clients", C_PURPLE, "Multi-device sync & identity switchers")
    bullet_block(tf2, [
        "Stateless clients: the Next.js web app (zkd-app/) and the Expo React Native app (zkd-android/).",
        "Clean polling architecture: GET /api/passengers/[id]/schedule.",
        "Identity switcher (?as=p-priya vs ?as=p-arjun): two tabs side by side prove independent tenant state.",
        "Actions on mobile and web instantly converge via the shared backend.",
    ], first_new=True)

    # ==================== SLIDE 8: Honest Experience KPIs Matrix ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "07", "Experience KPIs", "Customer experience KPIs",
                "Strictly no revenue proxies — every metric ties to a measured system evidence source.", 8)

    rows, cols = 5, 4
    left = Inches(0.5)
    top = Inches(1.65)
    width = Inches(12.33)
    height = Inches(4.9)

    table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.33)
    table.columns[2].width = Inches(1.9)
    table.columns[3].width = Inches(4.9)

    headers = ["Category", "Metric name", "Status", "System evidence / reality"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_CARD_DIM
        cell.margin_left = Inches(0.14)
        cell.margin_top = Inches(0.08)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.name = FONT_BODY
        p.font.color.rgb = C_ACCENT

    kpi_data = [
        ("Speed", "Time to Plan Ready (A2)", "EXISTS", "~11s measured via the pipeline journal (journal.ts)."),
        ("Member effort", "Hand-off Rate (B2)", "EXISTS", "Tracked via DisruptionResolution.kind === 'handed-over'."),
        ("Outcome quality", "Plan Acceptance Rate (C1)", "EXISTS", "Resolved as approved/autopilot over total recoveries."),
        ("Trust & loyalty", "Autopilot Opt-in (E2)", "EXISTS", "Share of members trusting the system unsupervised."),
    ]

    for i, row_vals in enumerate(kpi_data):
        for j, val in enumerate(row_vals):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD if i % 2 == 0 else C_BG_GLOW2
            cell.margin_left = Inches(0.14)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = FONT_BODY
            p.font.color.rgb = C_BODY
            if j == 2:
                p.font.bold = True
                p.font.color.rgb = C_GOOD
            if j == 0:
                p.font.color.rgb = C_WHITE

    # ==================== SLIDE 9: The 11-Second Recovery Breakdown ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "08", "Engineering Deep Dive", "The 11-second recovery",
                "Where the time actually goes, and why prediction — not raw speed — is the real unlock.", 9)

    box_rec1 = panel(slide9, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.15))
    tf1 = box_rec1.text_frame
    pad(tf1)
    panel_label(tf1, "Breakdown", C_ACCENT, "Where does the ~11s go?")
    bullet_block(tf1, [
        "95% of wall-clock time (~10.4s) is pure network I/O waiting on external supplier APIs (Duffel flight booking, LiteAPI hotel reservation).",
        "Under 0.6s is computational overhead: preference scoring, allocation, 3 negotiation rounds, and OPA policy checks.",
        "Zero GPU on the critical path — rule-based preference matching and heuristic ranking run instantly in-memory.",
    ], size=13, first_new=True)

    box_rec2 = panel(slide9, Inches(6.65), Inches(1.55), Inches(6.18), Inches(5.15))
    tf2 = box_rec2.text_frame
    pad(tf2)
    panel_label(tf2, "Why it matters", C_WARN, "Why prediction is the enabler")
    bullet_block(tf2, [
        "Prediction buys speed, not safety. A false positive costs one API call; a false negative costs 42 seconds.",
        "Pre-caching alternative itineraries before cancellation means the choice set is already in memory when consent arrives.",
        "Safety rests entirely on consent gates and notification ladders — never on a probability score.",
    ], size=13, first_new=True)

    # ==================== SLIDE 10: Limitations, Roadmap & Conclusion ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "09", "Limitations & Roadmap", "Honest limitations, roadmap & conclusion",
                "What is not built yet, and what ships next.", 10)

    box_lim1 = panel(slide10, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.15))
    tf1 = box_lim1.text_frame
    pad(tf1)
    panel_label(tf1, "Limitations", C_BAD, "Honest limitations")
    bullet_block(tf1, [
        "No Indian domestic training data yet: the risk model cold-starts to a population base rate until the retrain loop runs.",
        "Payment is mocked — a vPayment contract test, no live payment gateway.",
        "Pull-only API rate-limit constraints on flight status feeds.",
    ], marker_color=C_BAD, first_new=True)

    box_lim2 = panel(slide10, Inches(6.65), Inches(1.55), Inches(6.18), Inches(5.15))
    tf2 = box_lim2.text_frame
    pad(tf2)
    panel_label(tf2, "Roadmap", C_GOOD, "Roadmap & closing")
    bullet_block(tf2, [
        "Production rollout: live webhook subscriptions and automated event loggers.",
        "Full integration of the default-deny policy engine (server/policy/).",
        "Elevating American Express cardmember trust through uncompromising proactive service.",
    ], marker_color=C_GOOD, first_new=True)

    # Save presentation
    repo_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_path = os.path.join(repo_root, "assets", "deck", "ZKD_Concierge_Codestreet_2026.pptx")
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")


if __name__ == '__main__':
    create_deck()
